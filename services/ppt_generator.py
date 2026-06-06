import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from models.schemas import FullReport, Priority, Severity, Status

# ── Colour palette (matches PDF + Streamlit UI) ───────────────────────────────
C_DARK        = RGBColor(0x0a, 0x0a, 0x0f)
C_SURFACE     = RGBColor(0x11, 0x11, 0x18)
C_BORDER      = RGBColor(0x1e, 0x1e, 0x2e)
C_PURPLE      = RGBColor(0x6c, 0x63, 0xff)
C_PURPLE_SOFT = RGBColor(0xa7, 0x8b, 0xfa)
C_WHITE       = RGBColor(0xff, 0xff, 0xff)
C_TEXT_PRI    = RGBColor(0xe8, 0xe8, 0xf0)
C_TEXT_SEC    = RGBColor(0x90, 0x90, 0xa8)
C_GREEN       = RGBColor(0x22, 0xc5, 0x5e)
C_YELLOW      = RGBColor(0xea, 0xb3, 0x08)
C_RED         = RGBColor(0xef, 0x44, 0x44)
C_BLUE        = RGBColor(0x3b, 0x82, 0xf6)

PRIORITY_COLOR = {
    "HIGH":   C_RED,
    "MEDIUM": C_YELLOW,
    "LOW":    C_GREEN,
}
STATUS_COLOR = {
    "DECIDED":  C_GREEN,
    "DEFERRED": C_YELLOW,
    "PENDING":  C_RED,
}

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _rgb(r, g, b): return RGBColor(r, g, b)

def _bg(slide, color: RGBColor = None):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C_DARK

def _rect(slide, x, y, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _textbox(slide, x, y, w, h, text, size=12, bold=False,
             color: RGBColor = None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_TEXT_PRI
    return txBox

def _label(slide, x, y, w, h, text, bg: RGBColor, fg: RGBColor = None):
    """Coloured pill label (badge)."""
    _rect(slide, x, y, w, h, bg)
    _textbox(slide, x + Inches(0.05), y + Inches(0.02),
             w - Inches(0.1), h,
             text, size=7, bold=True,
             color=fg or C_WHITE, align=PP_ALIGN.CENTER)

def _divider(slide, y, color: RGBColor = None):
    line = slide.shapes.add_shape(1, Inches(0.4), y, W - Inches(0.8), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color or C_BORDER
    line.line.fill.background()

def _section_title(slide, y, text, emoji=""):
    label = f"{emoji}  {text}" if emoji else text
    _textbox(slide, Inches(0.4), y, W - Inches(0.8), Inches(0.4),
             label, size=16, bold=True, color=C_PURPLE)
    _divider(slide, y + Inches(0.38), C_PURPLE)


# ── Slide factories ────────────────────────────────────────────────────────────

def _slide_cover(prs: Presentation, report: FullReport):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    _bg(slide, C_DARK)

    # Top purple bar
    _rect(slide, 0, 0, W, Inches(0.08), C_PURPLE)

    # MeetingIQ wordmark
    _textbox(slide, Inches(0.5), Inches(0.25), Inches(5), Inches(0.6),
             "MeetingIQ", size=32, bold=True, color=C_WHITE)
    _textbox(slide, Inches(0.5), Inches(0.82), Inches(6), Inches(0.3),
             "AI Meeting Intelligence Platform", size=11, color=C_TEXT_SEC)

    # Generated date (top right)
    generated = datetime.now().strftime("%B %d, %Y  •  %H:%M")
    _textbox(slide, W - Inches(3.5), Inches(0.3), Inches(3.2), Inches(0.3),
             generated, size=9, color=C_TEXT_SEC, align=PP_ALIGN.RIGHT)

    # Meeting type badge
    mt = report.summary.meeting_type or "Meeting"
    _label(slide, Inches(0.5), Inches(1.5), Inches(1.8), Inches(0.32),
           mt.upper(), C_PURPLE)

    # One-line title
    _textbox(slide, Inches(0.5), Inches(1.95), W - Inches(1), Inches(1.2),
             report.summary.one_line,
             size=28, bold=True, color=C_WHITE)

    # Stat cards row
    stats = [
        ("ACTION ITEMS", str(len(report.action_items))),
        ("RISKS",        str(len(report.risks))),
        ("FINANCIALS",   str(len(report.financial_items))),
        ("PEOPLE",       str(len(report.people))),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.35)
    gap    = Inches(0.25)
    start_x = Inches(0.5)
    card_y  = Inches(5.4)

    for i, (label, val) in enumerate(stats):
        cx = start_x + i * (card_w + gap)
        _rect(slide, cx, card_y, card_w, card_h, C_SURFACE)
        _textbox(slide, cx + Inches(0.15), card_y + Inches(0.12),
                 card_w - Inches(0.3), Inches(0.28),
                 label, size=7, bold=True, color=C_TEXT_SEC)
        _textbox(slide, cx + Inches(0.15), card_y + Inches(0.38),
                 card_w - Inches(0.3), Inches(0.8),
                 val, size=40, bold=True, color=C_PURPLE)

    # Bottom purple bar
    _rect(slide, 0, H - Inches(0.18), W, Inches(0.18), C_PURPLE)


def _slide_summary(prs: Presentation, report: FullReport):
    s = report.summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _section_title(slide, Inches(0.2), "Core Summary & Decisions", "📋")

    # Decisions table header
    col_x = [Inches(0.4), Inches(7.5), Inches(10.2)]
    headers = ["DECISION", "STATUS", "OWNER"]
    header_y = Inches(0.85)
    col_w = [Inches(6.8), Inches(2.4), Inches(3.0)]

    for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        _rect(slide, cx, header_y, cw - Inches(0.05), Inches(0.28), C_SURFACE)
        _textbox(slide, cx + Inches(0.08), header_y + Inches(0.04),
                 cw - Inches(0.2), Inches(0.22),
                 hdr, size=7, bold=True, color=C_PURPLE_SOFT)

    row_y = header_y + Inches(0.32)
    row_h = Inches(0.38)
    for dec in s.decisions_made[:8]:
        bg = C_DARK if (s.decisions_made.index(dec) % 2 == 0) else C_SURFACE
        for cx, cw in zip(col_x, col_w):
            _rect(slide, cx, row_y, cw - Inches(0.05), row_h, bg)

        _textbox(slide, col_x[0] + Inches(0.08), row_y + Inches(0.06),
                 col_w[0] - Inches(0.2), row_h, dec.decision, size=9, color=C_TEXT_PRI)

        sc = STATUS_COLOR.get(dec.status.value, C_TEXT_SEC)
        _label(slide, col_x[1] + Inches(0.08), row_y + Inches(0.06),
               Inches(1.8), Inches(0.24), dec.status.value, sc)

        _textbox(slide, col_x[2] + Inches(0.08), row_y + Inches(0.06),
                 col_w[2] - Inches(0.2), row_h,
                 dec.owner or "—", size=9, color=C_TEXT_SEC)
        row_y += row_h + Inches(0.04)

    # Next steps (right side or below if space allows)
    ns_y = row_y + Inches(0.2)
    _textbox(slide, Inches(0.4), ns_y, Inches(4), Inches(0.28),
             "NEXT STEPS", size=8, bold=True, color=C_PURPLE_SOFT)
    ns_y += Inches(0.3)
    for i, step in enumerate(s.next_steps[:5], 1):
        _textbox(slide, Inches(0.4), ns_y, W - Inches(0.8), Inches(0.28),
                 f"{i}.  {step}", size=9, color=C_TEXT_PRI)
        ns_y += Inches(0.3)


def _slide_action_items(prs: Presentation, report: FullReport):
    if not report.action_items:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, Inches(0.2), "Action Items", "✅")

    col_x = [Inches(0.4), Inches(3.5), Inches(9.8), Inches(11.8)]
    col_w = [Inches(2.8), Inches(6.0), Inches(1.8), Inches(1.3)]
    headers = ["OWNER", "TASK", "DEADLINE", "PRIORITY"]

    header_y = Inches(0.85)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        _rect(slide, cx, header_y, cw - Inches(0.05), Inches(0.28), C_SURFACE)
        _textbox(slide, cx + Inches(0.08), header_y + Inches(0.04),
                 cw - Inches(0.2), Inches(0.22),
                 hdr, size=7, bold=True, color=C_PURPLE_SOFT)

    row_y = header_y + Inches(0.32)
    row_h = Inches(0.38)
    for i, item in enumerate(report.action_items[:12]):
        bg = C_DARK if i % 2 == 0 else C_SURFACE
        for cx, cw in zip(col_x, col_w):
            _rect(slide, cx, row_y, cw - Inches(0.05), row_h, bg)

        _textbox(slide, col_x[0] + Inches(0.08), row_y + Inches(0.06),
                 col_w[0] - Inches(0.2), row_h,
                 item.owner, size=9, bold=True, color=C_TEXT_PRI)
        _textbox(slide, col_x[1] + Inches(0.08), row_y + Inches(0.06),
                 col_w[1] - Inches(0.2), row_h,
                 item.task, size=9, color=C_TEXT_PRI)
        _textbox(slide, col_x[2] + Inches(0.08), row_y + Inches(0.06),
                 col_w[2] - Inches(0.2), row_h,
                 item.deadline or "—", size=8, color=C_TEXT_SEC)
        pc = PRIORITY_COLOR.get(item.priority.value, C_TEXT_SEC)
        _label(slide, col_x[3] + Inches(0.05), row_y + Inches(0.07),
               Inches(1.1), Inches(0.22), item.priority.value, pc)
        row_y += row_h + Inches(0.04)


def _slide_financials(prs: Presentation, report: FullReport):
    if not report.financial_items:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, Inches(0.2), "Financial Items", "💰")

    col_x = [Inches(0.4), Inches(2.6), Inches(9.0), Inches(11.2)]
    col_w = [Inches(2.0), Inches(6.2), Inches(2.0), Inches(1.9)]
    headers = ["AMOUNT", "CONTEXT", "STATUS", "OWNER"]

    header_y = Inches(0.85)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        _rect(slide, cx, header_y, cw - Inches(0.05), Inches(0.28), C_SURFACE)
        _textbox(slide, cx + Inches(0.08), header_y + Inches(0.04),
                 cw - Inches(0.2), Inches(0.22),
                 hdr, size=7, bold=True, color=C_PURPLE_SOFT)

    row_y = header_y + Inches(0.32)
    row_h = Inches(0.38)
    for i, item in enumerate(report.financial_items[:12]):
        bg = C_DARK if i % 2 == 0 else C_SURFACE
        for cx, cw in zip(col_x, col_w):
            _rect(slide, cx, row_y, cw - Inches(0.05), row_h, bg)

        _textbox(slide, col_x[0] + Inches(0.08), row_y + Inches(0.06),
                 col_w[0] - Inches(0.2), row_h,
                 item.amount, size=9, bold=True, color=C_PURPLE_SOFT)
        _textbox(slide, col_x[1] + Inches(0.08), row_y + Inches(0.06),
                 col_w[1] - Inches(0.2), row_h,
                 item.context, size=9, color=C_TEXT_PRI)
        sc = STATUS_COLOR.get(item.status.value, C_TEXT_SEC)
        _label(slide, col_x[2] + Inches(0.05), row_y + Inches(0.07),
               Inches(1.7), Inches(0.22), item.status.value, sc)
        _textbox(slide, col_x[3] + Inches(0.08), row_y + Inches(0.06),
                 col_w[3] - Inches(0.2), row_h,
                 item.owner or "—", size=8, color=C_TEXT_SEC)
        row_y += row_h + Inches(0.04)


def _slide_risks(prs: Presentation, report: FullReport):
    if not report.risks:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, Inches(0.2), "Risks & Blockers", "⚠️")

    col_x = [Inches(0.4), Inches(5.8), Inches(7.8), Inches(10.0)]
    col_w = [Inches(5.2), Inches(1.8), Inches(2.0), Inches(3.2)]
    headers = ["RISK", "SEVERITY", "OWNER", "MITIGATION"]

    header_y = Inches(0.85)
    for hdr, cx, cw in zip(headers, col_x, col_w):
        _rect(slide, cx, header_y, cw - Inches(0.05), Inches(0.28), C_SURFACE)
        _textbox(slide, cx + Inches(0.08), header_y + Inches(0.04),
                 cw - Inches(0.2), Inches(0.22),
                 hdr, size=7, bold=True, color=C_PURPLE_SOFT)

    row_y = header_y + Inches(0.32)
    row_h = Inches(0.42)
    for i, item in enumerate(report.risks[:10]):
        bg = C_DARK if i % 2 == 0 else C_SURFACE
        for cx, cw in zip(col_x, col_w):
            _rect(slide, cx, row_y, cw - Inches(0.05), row_h, bg)

        _textbox(slide, col_x[0] + Inches(0.08), row_y + Inches(0.06),
                 col_w[0] - Inches(0.2), row_h,
                 item.description, size=9, color=C_TEXT_PRI)
        sc = PRIORITY_COLOR.get(item.severity.value, C_TEXT_SEC)
        _label(slide, col_x[1] + Inches(0.05), row_y + Inches(0.09),
               Inches(1.5), Inches(0.22), item.severity.value, sc)
        _textbox(slide, col_x[2] + Inches(0.08), row_y + Inches(0.06),
                 col_w[2] - Inches(0.2), row_h,
                 item.owner or "—", size=8, color=C_TEXT_SEC)
        _textbox(slide, col_x[3] + Inches(0.08), row_y + Inches(0.06),
                 col_w[3] - Inches(0.2), row_h,
                 item.mitigation or "—", size=8, color=C_TEXT_PRI)
        row_y += row_h + Inches(0.04)


def _slide_people_dates(prs: Presentation, report: FullReport):
    """People & Roles + Dates & Deadlines on one slide (two columns)."""
    if not report.people and not report.dates_deadlines:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    # Left — People
    if report.people:
        _textbox(slide, Inches(0.4), Inches(0.2), Inches(5.5), Inches(0.38),
                 "👥  People & Roles", size=16, bold=True, color=C_PURPLE)
        _divider(slide, Inches(0.58), C_PURPLE)
        row_y = Inches(0.75)
        for i, p in enumerate(report.people[:10]):
            bg = C_DARK if i % 2 == 0 else C_SURFACE
            _rect(slide, Inches(0.4), row_y, Inches(6.0), Inches(0.38), bg)
            _textbox(slide, Inches(0.5), row_y + Inches(0.05),
                     Inches(2.2), Inches(0.28),
                     p.name, size=9, bold=True, color=C_TEXT_PRI)
            _textbox(slide, Inches(2.75), row_y + Inches(0.05),
                     Inches(2.0), Inches(0.28),
                     p.role or "—", size=8, color=C_TEXT_SEC)
            _textbox(slide, Inches(4.8), row_y + Inches(0.05),
                     Inches(1.5), Inches(0.28),
                     p.email or "", size=7, color=C_PURPLE_SOFT)
            row_y += Inches(0.42)

    # Right — Dates
    if report.dates_deadlines:
        _textbox(slide, Inches(7.0), Inches(0.2), Inches(5.5), Inches(0.38),
                 "📅  Dates & Deadlines", size=16, bold=True, color=C_PURPLE)
        _divider(slide, Inches(0.58), C_PURPLE)
        row_y = Inches(0.75)
        for i, d in enumerate(report.dates_deadlines[:10]):
            bg = C_DARK if i % 2 == 0 else C_SURFACE
            _rect(slide, Inches(7.0), row_y, Inches(6.1), Inches(0.38), bg)
            _textbox(slide, Inches(7.1), row_y + Inches(0.05),
                     Inches(1.6), Inches(0.28),
                     d.date, size=9, bold=True, color=C_PURPLE_SOFT)
            _textbox(slide, Inches(8.75), row_y + Inches(0.05),
                     Inches(3.2), Inches(0.28),
                     d.event, size=8, color=C_TEXT_PRI)
            pc = PRIORITY_COLOR.get(d.priority.value, C_TEXT_SEC)
            _label(slide, Inches(12.0), row_y + Inches(0.07),
                   Inches(0.9), Inches(0.22), d.priority.value[0], pc)
            row_y += Inches(0.42)


def _slide_questions_commitments(prs: Presentation, report: FullReport):
    """Unresolved Questions + Commitments on one slide."""
    if not report.unresolved_questions and not report.commitments:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    if report.unresolved_questions:
        _textbox(slide, Inches(0.4), Inches(0.2), Inches(6), Inches(0.38),
                 "❓  Unresolved Questions", size=16, bold=True, color=C_PURPLE)
        _divider(slide, Inches(0.58), C_PURPLE)
        row_y = Inches(0.75)
        for i, q in enumerate(report.unresolved_questions[:7]):
            bg = C_DARK if i % 2 == 0 else C_SURFACE
            _rect(slide, Inches(0.4), row_y, Inches(6.2), Inches(0.5), bg)
            _textbox(slide, Inches(0.5), row_y + Inches(0.04),
                     Inches(6.0), Inches(0.25),
                     q.question, size=9, bold=True, color=C_TEXT_PRI)
            _textbox(slide, Inches(0.5), row_y + Inches(0.26),
                     Inches(6.0), Inches(0.22),
                     q.context, size=7, color=C_TEXT_SEC)
            row_y += Inches(0.54)

    if report.commitments:
        _textbox(slide, Inches(7.0), Inches(0.2), Inches(6), Inches(0.38),
                 "🤝  Commitments", size=16, bold=True, color=C_PURPLE)
        _divider(slide, Inches(0.58), C_PURPLE)
        row_y = Inches(0.75)
        for i, c in enumerate(report.commitments[:7]):
            bg = C_DARK if i % 2 == 0 else C_SURFACE
            _rect(slide, Inches(7.0), row_y, Inches(6.1), Inches(0.5), bg)
            _textbox(slide, Inches(7.1), row_y + Inches(0.04),
                     Inches(2.0), Inches(0.25),
                     c.person, size=9, bold=True, color=C_PURPLE_SOFT)
            _textbox(slide, Inches(9.1), row_y + Inches(0.04),
                     Inches(3.8), Inches(0.25),
                     c.committed_to, size=9, color=C_TEXT_PRI)
            _textbox(slide, Inches(7.1), row_y + Inches(0.27),
                     Inches(6.0), Inches(0.2),
                     f"Target: {c.target or '—'}   •   By: {c.deadline or '—'}",
                     size=7, color=C_TEXT_SEC)
            row_y += Inches(0.54)


def _slide_closing(prs: Presentation, report: FullReport):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_DARK)

    _rect(slide, 0, 0, W, Inches(0.08), C_PURPLE)

    _textbox(slide, Inches(0.5), Inches(2.4), W - Inches(1), Inches(1.0),
             "MeetingIQ", size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.5), Inches(3.4), W - Inches(1), Inches(0.4),
             "AI Meeting Intelligence", size=14, color=C_TEXT_SEC, align=PP_ALIGN.CENTER)

    generated = datetime.now().strftime("%B %d, %Y")
    _textbox(slide, Inches(0.5), Inches(4.2), W - Inches(1), Inches(0.35),
             f"Report generated {generated}", size=10,
             color=C_TEXT_SEC, align=PP_ALIGN.CENTER, italic=True)

    _rect(slide, 0, H - Inches(0.18), W, Inches(0.18), C_PURPLE)


# ── Public entry point ────────────────────────────────────────────────────────

def generate_ppt(report: FullReport, output_path: str) -> str:
    """
    Generates a styled .pptx from a FullReport object.
    Returns output_path on success.
    """
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _slide_cover(prs, report)
    _slide_summary(prs, report)
    _slide_action_items(prs, report)
    _slide_financials(prs, report)
    _slide_risks(prs, report)
    _slide_people_dates(prs, report)
    _slide_questions_commitments(prs, report)
    _slide_closing(prs, report)

    prs.save(output_path)
    print(f"[PPT] Generated → {output_path}")
    return output_path